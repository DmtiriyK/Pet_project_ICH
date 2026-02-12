from __future__ import annotations

import json
from dataclasses import dataclass
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import base64


ROOT = Path(__file__).resolve().parents[1]
FINAL_DIR = ROOT / "reports" / "final"
EDA_FIG = ROOT / "reports" / "eda" / "figures"
TIME_FIG = ROOT / "reports" / "time" / "figures"


@dataclass(frozen=True)
class FigureSlide:
    title: str
    path: Path
    caption: str | None = None


def _read_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def _fmt_money(x: float) -> str:
    return f"{x:,.0f}"


def _fmt_pct(x: float) -> str:
    return f"{x:.2%}"


def _ppt_add_title(slide, text: str) -> None:
    title = slide.shapes.title
    title.text = text
    for p in title.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(36)


def _ppt_add_bullets(slide, bullets: list[str]) -> None:
    left, top, width, height = Inches(1.0), Inches(1.7), Inches(11.3), Inches(5.2)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(22)


def _ppt_add_caption(slide, caption: str) -> None:
    left, top, width, height = Inches(1.0), Inches(6.9), Inches(11.3), Inches(0.6)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.text = caption
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.LEFT


def _ppt_add_image(slide, img: Path) -> None:
    left, top, width, height = Inches(0.8), Inches(1.6), Inches(12.0), Inches(5.2)
    slide.shapes.add_picture(str(img), left, top, width=width, height=height)


def build_pptx() -> Path:
    FINAL_DIR.mkdir(parents=True, exist_ok=True)
    out_path = FINAL_DIR / "presentation.pptx"

    overall = _read_json(ROOT / "reports" / "eda" / "metrics_overall.json")
    ttc = _read_json(ROOT / "reports" / "time" / "paid_time_to_close_stats.json")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "CRM аналитика — итоговый проект"
    subtitle = slide.placeholders[1]
    subtitle.text = "Главный инсайт №1: продажи (разброс по менеджерам)"

    # Agenda
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _ppt_add_title(slide, "План")
    _ppt_add_bullets(
        slide,
        [
            "Данные и правила (что считаем оплатой)",
            "Воронка и ключевые метрики (cash + contract)",
            "Продажи — главный инсайт",
            "Реклама, продукты, сегменты",
            "Время (time-to-close) и гипотезы на 2 недели",
        ],
    )

    # Key numbers
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _ppt_add_title(slide, "Ключевые числа (окно пересечения Spend и Deals)")
    _ppt_add_bullets(
        slide,
        [
            f"Spend: {_fmt_money(overall['spend_total'])}",
            f"Deals: {int(overall['deals_total']):,}",
            f"Paid deals: {int(overall['paid_deals']):,} (paid rate {_fmt_pct(overall['paid_rate'])})",
            f"Revenue cash: {_fmt_money(overall['revenue_cash_total'])}",
            f"Revenue contract: {_fmt_money(overall['revenue_contract_total'])}",
        ],
    )

    # Funnel
    funnel_img = EDA_FIG / "stage_funnel_top12.png"
    if funnel_img.exists():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        _ppt_add_title(slide, "Воронка (Stage)")
        _ppt_add_image(slide, funnel_img)
        _ppt_add_caption(slide, "Источник: reports/eda/figures/stage_funnel_top12.png")

    # Sales insight
    sales_img = EDA_FIG / "paid_rate_by_owner_top15.png"
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _ppt_add_title(slide, "Продажи (инсайт №1): разброс по менеджерам")
    if sales_img.exists():
        _ppt_add_image(slide, sales_img)
        _ppt_add_caption(slide, "Видно сильный разброс paid_rate по deal_owner_name.")
    else:
        _ppt_add_bullets(slide, ["Сильный разброс paid_rate по менеджерам → конверсия зависит от процесса/исполнителя."])

    # Sales hypothesis
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _ppt_add_title(slide, "Гипотеза роста (продажи) — 2 недели")
    _ppt_add_bullets(
        slide,
        [
            "Если реплицировать практики топ-менеджеров (скрипт, SLA, квалификация), то paid_rate вырастет.",
            "Метрики: paid_rate и revenue_contract по deal_owner_name.",
            "Дизайн: пилотная группа A vs контроль B, одинаковые источники лидов.",
            "Критерий успеха: рост paid_rate в A при сохранении объёма обработанных сделок.",
        ],
    )

    # Ads
    ads_img = EDA_FIG / "contract_roas_by_source.png"
    if ads_img.exists():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        _ppt_add_title(slide, "Реклама: ROAS по Source (контракт)")
        _ppt_add_image(slide, ads_img)
        _ppt_add_caption(slide, "Сравнение источников по contract ROAS (с фильтром spend/paid).")

    # Products
    prod_img = EDA_FIG / "revenue_contract_by_product_top15.png"
    if prod_img.exists():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        _ppt_add_title(slide, "Продукты: выручка (контракт) по paid-only")
        _ppt_add_image(slide, prod_img)
        _ppt_add_caption(slide, "paid-only, чтобы сравнивать AOV/выручку по продуктам.")

    # Time-to-close
    ttc_img = TIME_FIG / "time_to_close_hist.png"
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _ppt_add_title(slide, "Время до оплаты (time-to-close)")
    if ttc_img.exists():
        _ppt_add_image(slide, ttc_img)
        _ppt_add_caption(
            slide,
            f"Coverage closing_date у paid: {ttc['coverage_pct']:.2f}%. Median lag: {ttc['lag_days_median']:.2f} дней.",
        )
    else:
        _ppt_add_bullets(
            slide,
            [
                f"Coverage closing_date у paid: {ttc['coverage_pct']:.2f}%",
                f"Median lag: {ttc['lag_days_median']:.2f} дней; P90: {ttc['lag_days_p90']:.2f} дней",
            ],
        )

    # Caveats
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _ppt_add_title(slide, "Оговорки по данным")
    _ppt_add_bullets(
        slide,
        [
            "Оплата: только Stage = Payment Done.",
            "Closing Date у части paid пустой → time-to-close только по подвыборке.",
            "ID контактов в Calls/Deals пришли как Excel-числа → точный джойн Contacts↔Calls↔Deals не гарантируется.",
            "В отчёте держим две выручки: cash (Initial) и contract (Offer Total).",
        ],
    )

    prs.save(out_path)
    return out_path


def _img_data_uri(path: Path) -> str:
    data = path.read_bytes()
    b64 = base64.b64encode(data).decode("ascii")
    return f"data:image/png;base64,{b64}"


def build_html() -> Path:
    FINAL_DIR.mkdir(parents=True, exist_ok=True)
    out_path = FINAL_DIR / "slides.html"

    overall = _read_json(ROOT / "reports" / "eda" / "metrics_overall.json")
    ttc = _read_json(ROOT / "reports" / "time" / "paid_time_to_close_stats.json")

    def uri(p: Path) -> str:
        # embed image as data URI to guarantee it appears in PDF
        return _img_data_uri(p)

    slides = [
        (
            "CRM аналитика — итоговый проект",
            "<p>Главный инсайт №1: <b>продажи</b> (разброс по менеджерам)</p>",
        ),
        (
            "Ключевые числа",
            f"""
<ul>
  <li>Spend: <b>{_fmt_money(overall['spend_total'])}</b></li>
  <li>Deals: <b>{int(overall['deals_total']):,}</b></li>
  <li>Paid deals: <b>{int(overall['paid_deals']):,}</b> (paid rate <b>{_fmt_pct(overall['paid_rate'])}</b>)</li>
  <li>Revenue cash: <b>{_fmt_money(overall['revenue_cash_total'])}</b></li>
  <li>Revenue contract: <b>{_fmt_money(overall['revenue_contract_total'])}</b></li>
</ul>
""",
        ),
        (
            "Продажи (инсайт №1)",
            f"<img class='fit' src='{uri(EDA_FIG / 'paid_rate_by_owner_top15.png')}' />",
        ),
        (
            "Гипотеза роста (продажи) — 2 недели",
            """
<ul>
  <li>Реплицировать практики топ-менеджеров (скрипт, SLA, квалификация).</li>
  <li>Метрики: paid_rate и revenue_contract по deal_owner_name.</li>
  <li>Дизайн: группа A vs контроль B, одинаковые источники лидов.</li>
  <li>Успех: рост paid_rate в A без падения объёма.</li>
</ul>
""",
        ),
        ("Воронка (Stage)", f"<img class='fit' src='{uri(EDA_FIG / 'stage_funnel_top12.png')}' />"),
        ("Реклама: ROAS по Source (контракт)", f"<img class='fit' src='{uri(EDA_FIG / 'contract_roas_by_source.png')}' />"),
        ("Продукты: выручка (контракт) paid-only", f"<img class='fit' src='{uri(EDA_FIG / 'revenue_contract_by_product_top15.png')}' />"),
        (
            "Время до оплаты (time-to-close)",
            f"""
<p>Coverage closing_date у paid: <b>{ttc['coverage_pct']:.2f}%</b>. Median lag: <b>{ttc['lag_days_median']:.2f}</b> дней.</p>
<img class='fit' src='{uri(TIME_FIG / 'time_to_close_hist.png')}' />
""",
        ),
        (
            "Оговорки по данным",
            """
<ul>
  <li>Оплата: только Stage = Payment Done.</li>
  <li>Closing Date у части paid пустой → time-to-close по подвыборке.</li>
  <li>ID контактов в Calls/Deals — Excel-числа → точный джойн не гарантируется.</li>
  <li>Две выручки: cash (Initial) и contract (Offer Total).</li>
</ul>
""",
        ),
    ]

    sections_html = "\n".join(
        f"<section><h2>{title}</h2>{body}</section>"
        if title != "CRM аналитика — итоговый проект"
        else f"<section><h1>{title}</h1>{body}</section>"
        for title, body in slides
    )

    html = f"""<!doctype html>
<html lang="ru">
  <head>
    <meta charset="utf-8" />
    <meta name="viewport" content="width=device-width, initial-scale=1.0" />
    <title>CRM аналитика — презентация</title>
    <link rel="stylesheet" href="https://unpkg.com/reveal.js/dist/reveal.css" />
    <link rel="stylesheet" href="https://unpkg.com/reveal.js/dist/theme/white.css" />
    <style>
      .reveal .slides {{ text-align: left; }}
      .reveal h1, .reveal h2 {{ text-align: left; }}
      img.fit {{
        max-width: 100%;
        max-height: 520px;
        width: auto;
        height: auto;
        border: 1px solid #eee;
      }}
    </style>
  </head>
  <body>
    <div class="reveal">
      <div class="slides">
        {sections_html}
      </div>
    </div>
    <script src="https://unpkg.com/reveal.js/dist/reveal.js"></script>
    <script>
      Reveal.initialize({{
        hash: true,
        slideNumber: true,
      }});
    </script>
  </body>
</html>
"""

    out_path.write_text(html, encoding="utf-8")
    return out_path


def main() -> None:
    build_pptx()
    build_html()
    print("OK: reports/final/presentation.pptx")
    print("OK: reports/final/slides.html")


if __name__ == "__main__":
    main()
